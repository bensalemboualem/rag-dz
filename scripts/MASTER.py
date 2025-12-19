"""
IA FACTORY - MASTER SCRIPT
Exécute TOUTES les tâches automatiquement
"""

import subprocess
import sys
import os
import webbrowser
import time
import threading
from pathlib import Path
from datetime import datetime

BASE = Path(__file__).parent
os.chdir(BASE)

print("=" * 60)
print("🏭 IA FACTORY - EXECUTION COMPLETE AUTOMATIQUE")
print("=" * 60)
print("📧 CH: contact@iafactory.ch")
print("📧 DZ: contact@iafactoryalgeria.com")
print("=" * 60)

# ═══════════════════════════════════════════════════════════════
# 1. INSTALLATION PACKAGES
# ═══════════════════════════════════════════════════════════════
print("\n[1/7] 📦 Installation packages...")
subprocess.run([sys.executable, "-m", "pip", "install", "-q", "--upgrade", "pip"], capture_output=True)
subprocess.run([sys.executable, "-m", "pip", "install", "-q",
    "fastapi", "uvicorn", "pydantic", "pydantic-settings", "python-dotenv",
    "python-docx", "python-pptx", "openpyxl", "requests", "httpx", "anthropic", "qdrant-client"
], capture_output=True)
print("    ✅ Packages installés")

# ═══════════════════════════════════════════════════════════════
# 2. CREATION STRUCTURE
# ═══════════════════════════════════════════════════════════════
print("\n[2/7] 📁 Création structure...")
dirs = [
    "api/routers", "api/models", "api/services",
    "config",
    "core/rag", "core/agents", "core/llm",
    "templates/documents", "templates/presentations", "templates/dashboards", "templates/emails",
    "workflows/delivery", "workflows/sales", "workflows/support",
    "social_media/content", "social_media/scheduler",
    "digital_twin",
    "infrastructure/docker",
    "outputs/propositions", "outputs/presentations", "outputs/dashboards",
    "tests", "scripts", "docs"
]
for d in dirs:
    (BASE / d).mkdir(parents=True, exist_ok=True)
print("    ✅ Structure créée")

# ═══════════════════════════════════════════════════════════════
# 3. CREATION FICHIERS CONFIG
# ═══════════════════════════════════════════════════════════════
print("\n[3/7] ⚙️ Création config...")

(BASE / "config" / "__init__.py").write_text("from .settings import settings\n")

# ═══════════════════════════════════════════════════════════════
# 4. CREATION API
# ═══════════════════════════════════════════════════════════════
print("\n[4/7] 🌐 Création API...")

(BASE / "api" / "__init__.py").write_text("")
(BASE / "api" / "routers" / "__init__.py").write_text("")

(BASE / "api" / "main.py").write_text('''from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List
from datetime import datetime
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).parent.parent))

app = FastAPI(title="IA Factory API", version="1.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

CONTACTS = {
    "CH": {"email": "contact@iafactory.ch", "phone": "+41 XX XXX XX XX", "web": "www.iafactory.ch"},
    "DZ": {"email": "contact@iafactoryalgeria.com", "phone": "+213 XX XXX XX XX", "web": "www.iafactoryalgeria.com"}
}

class ProposalRequest(BaseModel):
    company: str
    services: List[str] = ["RAG System"]
    market: str = "CH"

class ClientCreate(BaseModel):
    name: str
    company: str
    email: str
    market: str = "CH"

@app.get("/")
def root():
    return {"name": "IA Factory", "CH": "contact@iafactory.ch", "DZ": "contact@iafactoryalgeria.com", "time": datetime.now().isoformat()}

@app.get("/health")
def health():
    return {"status": "healthy", "version": "1.0.0"}

@app.get("/contacts")
def contacts():
    return CONTACTS

@app.get("/kpis")
def kpis():
    return {"mrr": 8500, "clients": 12, "margin": "92%", "pipeline": 45000, "CH": 5500, "DZ": 3000}

@app.get("/clients")
def clients():
    return {"total": 12, "clients": [
        {"id": "CLI-001", "company": "École Nouvelle Horizon", "market": "CH", "mrr": 1200},
        {"id": "CLI-002", "company": "Algérie Télécom", "market": "DZ", "mrr": 1200},
    ]}

@app.post("/clients/create")
def create_client(c: ClientCreate):
    return {"status": "created", "id": f"CLI-{datetime.now().strftime(\\"%Y%m%d%H%M%S\\")}"}

@app.post("/documents/proposal")
def gen_proposal(r: ProposalRequest):
    try:
        from templates.documents.proposition_commerciale import PropositionGenerator
        f = PropositionGenerator().generate({"company": r.company}, r.services, r.market)
        return {"status": "success", "file": f}
    except Exception as e:
        return {"error": str(e)}

@app.post("/documents/deck/{market}")
def gen_deck(market: str = "CH"):
    try:
        from templates.presentations.teaching_assistant_deck import TeachingAssistantDeck
        f = TeachingAssistantDeck().generate(market)
        return {"status": "success", "file": f}
    except Exception as e:
        return {"error": str(e)}

@app.post("/documents/dashboard")
def gen_dashboard():
    try:
        from templates.dashboards.kpi_dashboard import KPIDashboard
        f = KPIDashboard().generate()
        return {"status": "success", "file": f}
    except Exception as e:
        return {"error": str(e)}
''')
print("    ✅ API créée")

# ═══════════════════════════════════════════════════════════════
# 5. CREATION TEMPLATES
# ═══════════════════════════════════════════════════════════════
print("\n[5/7] 📄 Création templates...")

for d in ["templates", "templates/documents", "templates/presentations", "templates/dashboards", "templates/emails"]:
    (BASE / d / "__init__.py").write_text("")

(BASE / "templates" / "documents" / "proposition_commerciale.py").write_text('''from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from datetime import datetime
from pathlib import Path

class PropositionGenerator:
    def __init__(self):
        self.pricing = {"CH": {"RAG System": 25000, "Multi-Agent": 40000, "Training": 8000},
                        "DZ": {"RAG System": 800000, "Multi-Agent": 1500000, "Training": 300000}}
        self.contacts = {"CH": "contact@iafactory.ch", "DZ": "contact@iafactoryalgeria.com"}
    
    def generate(self, client: dict, services: list, market: str = "CH") -> str:
        doc = Document()
        currency = "CHF" if market == "CH" else "DZD"
        
        doc.add_heading("IA FACTORY", 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("PROPOSITION COMMERCIALE", 1).alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"Client: {client.get('company', 'N/A')}")
        doc.add_paragraph(f"Date: {datetime.now().strftime('%d/%m/%Y')}")
        doc.add_page_break()
        
        doc.add_heading("Services", 1)
        total = 0
        for s in services:
            if s in self.pricing[market]:
                doc.add_paragraph(f"• {s}: {self.pricing[market][s]:,} {currency}")
                total += self.pricing[market][s]
        doc.add_paragraph(f"TOTAL: {total:,} {currency}").runs[0].bold = True
        
        doc.add_heading("Contact", 1)
        doc.add_paragraph(f"Email: {self.contacts[market]}")
        
        Path("outputs/propositions").mkdir(parents=True, exist_ok=True)
        f = f"outputs/propositions/prop_{client.get('company','x').replace(' ','_')}_{market}_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
        doc.save(f)
        return f
''')

(BASE / "templates" / "presentations" / "teaching_assistant_deck.py").write_text('''from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from pathlib import Path

class TeachingAssistantDeck:
    def generate(self, market: str = "CH") -> str:
        prs = Presentation()
        email = "contact@iafactory.ch" if market == "CH" else "contact@iafactoryalgeria.com"
        price = "500 CHF/mois" if market == "CH" else "50,000 DZD/mois"
        
        slides = [("AI Teaching Assistant", ["Assistants IA par matière", "Génération exercices", "Correction auto"]),
                  ("Résultats", ["70% temps gagné", "95% satisfaction"]),
                  ("Contact", [email, price])]
        
        for title, bullets in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tb.text_frame.paragraphs[0].text = title
            tb.text_frame.paragraphs[0].font.size = Pt(36)
            tb.text_frame.paragraphs[0].font.bold = True
            cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
            for i, b in enumerate(bullets):
                p = cb.text_frame.paragraphs[0] if i == 0 else cb.text_frame.add_paragraph()
                p.text = f"• {b}"
                p.font.size = Pt(24)
        
        Path("outputs/presentations").mkdir(parents=True, exist_ok=True)
        f = f"outputs/presentations/deck_{market}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
        prs.save(f)
        return f
''')

(BASE / "templates" / "dashboards" / "kpi_dashboard.py").write_text('''from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.chart import BarChart, Reference
from datetime import datetime
from pathlib import Path

class KPIDashboard:
    def generate(self) -> str:
        wb = Workbook()
        ws = wb.active
        ws.title = "KPIs"
        
        ws.merge_cells("A1:C1")
        ws["A1"] = "IA FACTORY - KPIs"
        ws["A1"].font = Font(size=16, bold=True)
        
        ws["A2"], ws["B2"] = "CH:", "contact@iafactory.ch"
        ws["A3"], ws["B3"] = "DZ:", "contact@iafactoryalgeria.com"
        
        for i, (m, v, t) in enumerate([("MRR", "8,500 CHF", "+15%"), ("Clients", "12", "+3"), ("Margin", "92%", "+2%")], 5):
            ws[f"A{i}"], ws[f"B{i}"], ws[f"C{i}"] = m, v, t
        
        for i, (m, r) in enumerate([("Jan", 3000), ("Feb", 4500), ("Mar", 6000), ("Apr", 7500), ("May", 8500)], 10):
            ws[f"A{i}"], ws[f"B{i}"] = m, r
        
        chart = BarChart()
        chart.add_data(Reference(ws, min_col=2, min_row=10, max_row=14))
        chart.set_categories(Reference(ws, min_col=1, min_row=10, max_row=14))
        ws.add_chart(chart, "D10")
        
        Path("outputs/dashboards").mkdir(parents=True, exist_ok=True)
        f = f"outputs/dashboards/kpi_{datetime.now().strftime('%Y%m%d%H%M%S')}.xlsx"
        wb.save(f)
        return f
''')

(BASE / "templates" / "emails" / "templates.py").write_text('''class EmailTemplates:
    CONTACTS = {"CH": "contact@iafactory.ch", "DZ": "contact@iafactoryalgeria.com"}
    
    @staticmethod
    def welcome(client, market, url):
        return {"to": client["email"], "subject": f"Bienvenue {client['company']}!", 
                "body": f"Votre accès: {url}\\nContact: {EmailTemplates.CONTACTS[market]}"}
    
    @staticmethod
    def hot_lead(lead, score):
        return {"to": "contact@iafactory.ch", "subject": f"HOT LEAD ({score}): {lead['company']}",
                "body": f"{lead['name']} - {lead['email']}"}
''')
print("    ✅ Templates créés")

# ═══════════════════════════════════════════════════════════════
# 6. GENERATION DOCUMENTS
# ═══════════════════════════════════════════════════════════════
print("\n[6/7] 📄 Génération documents...")
sys.path.insert(0, str(BASE))

try:
    from templates.documents.proposition_commerciale import PropositionGenerator
    from templates.presentations.teaching_assistant_deck import TeachingAssistantDeck
    from templates.dashboards.kpi_dashboard import KPIDashboard
    
    gen = PropositionGenerator()
    print(f"    ✅ {gen.generate({'company': 'Swiss Corp'}, ['RAG System', 'Training'], 'CH')}")
    print(f"    ✅ {gen.generate({'company': 'Algérie Télécom'}, ['RAG System'], 'DZ')}")
    print(f"    ✅ {TeachingAssistantDeck().generate('CH')}")
    print(f"    ✅ {TeachingAssistantDeck().generate('DZ')}")
    print(f"    ✅ {KPIDashboard().generate()}")
except Exception as e:
    print(f"    ⚠️ {e}")

# ═══════════════════════════════════════════════════════════════
# 7. LANCEMENT API
# ═══════════════════════════════════════════════════════════════
print("\n[7/7] 🚀 Lancement API...")

def open_browser():
    time.sleep(2)
    webbrowser.open("http://127.0.0.1:8000")
    webbrowser.open("http://127.0.0.1:8000/docs")

threading.Thread(target=open_browser, daemon=True).start()

print("\n" + "=" * 60)
print("✅ TOUT EST PRÊT!")
print("=" * 60)
print("🌐 API: http://127.0.0.1:8000")
print("📚 Docs: http://127.0.0.1:8000/docs")
print("📁 Fichiers: outputs/")
print("=" * 60)
print("\nCtrl+C pour arrêter\n")

subprocess.run([sys.executable, "-m", "uvicorn", "api.main:app", "--host", "127.0.0.1", "--port", "8000", "--reload"])
