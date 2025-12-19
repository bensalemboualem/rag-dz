from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from pathlib import Path

class TeachingAssistantDeck:
    def generate(self, market: str = "CH") -> str:
        prs = Presentation()
        email = "contact@iafactory.ch" if market == "CH" else "contact@iafactoryalgeria.com"
        price = "500 CHF/mois" if market == "CH" else "50,000 DZD/mois"
        
        slides = [("AI Teaching Assistant", ["Assistants IA par matière", "Génération exercices", "Correction auto"]),
                  ("Résultats", ["70% temps gagné", "95% satisfaction"]),
                  ("Contact", [email, price])]
        
        for title, bullets in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tb.text_frame.paragraphs[0].text = title
            tb.text_frame.paragraphs[0].font.size = Pt(36)
            tb.text_frame.paragraphs[0].font.bold = True
            cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
            for i, b in enumerate(bullets):
                p = cb.text_frame.paragraphs[0] if i == 0 else cb.text_frame.add_paragraph()
                p.text = f"• {b}"
                p.font.size = Pt(24)
        
        Path("outputs/presentations").mkdir(parents=True, exist_ok=True)
        f = f"outputs/presentations/deck_{market}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
        prs.save(f)
        return f
