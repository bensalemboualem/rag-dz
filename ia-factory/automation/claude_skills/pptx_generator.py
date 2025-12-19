"""
IA Factory - Claude Skill: PPTX Generator
Génération de présentations PowerPoint professionnelles
"""

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any
from datetime import datetime
from enum import Enum
import os

router = APIRouter(prefix="/skills/pptx", tags=["Claude Skills - PPTX"])


class PresentationType(str, Enum):
    PITCH_DECK = "pitch_deck"           # Pitch investisseurs
    SALES_DECK = "sales_deck"           # Présentation commerciale
    COMPANY_INTRO = "company_intro"     # Présentation entreprise
    PRODUCT_DEMO = "product_demo"       # Démo produit
    TRAINING = "training"               # Formation
    REPORT = "report"                   # Rapport
    WORKSHOP = "workshop"               # Atelier
    WEBINAR = "webinar"                 # Webinaire


class SlideLayout(str, Enum):
    TITLE = "title"
    TITLE_CONTENT = "title_content"
    TWO_COLUMNS = "two_columns"
    COMPARISON = "comparison"
    QUOTE = "quote"
    IMAGE_FULL = "image_full"
    CHART = "chart"
    BULLETS = "bullets"
    SECTION = "section"
    THANK_YOU = "thank_you"


class ColorTheme(str, Enum):
    IA_FACTORY = "ia_factory"       # Bleu IA Factory
    ALGERIA = "algeria"             # Vert/Blanc/Rouge
    SWISS = "swiss"                 # Rouge/Blanc
    TECH = "tech"                   # Bleu foncé/Cyan
    NATURE = "nature"               # Vert
    DARK = "dark"                   # Noir/Doré
    LIGHT = "light"                 # Blanc/Gris


class SlideRequest(BaseModel):
    """Définition d'une slide"""
    layout: SlideLayout
    title: str
    subtitle: Optional[str] = None
    content: Any = None  # Peut être str, list, dict selon le layout
    notes: Optional[str] = None  # Notes du présentateur
    image_path: Optional[str] = None


class PresentationRequest(BaseModel):
    """Requête de génération de présentation"""
    presentation_type: PresentationType
    title: str
    subtitle: Optional[str] = None
    author: str = "Boualem Chebaki"
    company: str = "IA Factory"
    theme: ColorTheme = ColorTheme.IA_FACTORY
    slides: List[SlideRequest] = Field(default_factory=list)
    include_agenda: bool = True
    include_contact: bool = True
    language: str = "fr"


class GeneratedPresentation(BaseModel):
    """Présentation générée"""
    id: str
    filename: str
    filepath: str
    presentation_type: PresentationType
    slide_count: int
    created_at: datetime
    download_url: str


class PptxGenerator:
    """
    Générateur de présentations PPTX
    Utilise python-pptx pour créer des présentations professionnelles
    """
    
    # Palettes de couleurs
    THEMES = {
        ColorTheme.IA_FACTORY: {
            "primary": "1F4E79",      # Bleu foncé
            "secondary": "2E75B6",    # Bleu moyen
            "accent": "00B0F0",       # Bleu vif
            "background": "FFFFFF",
            "text": "333333",
            "text_light": "FFFFFF"
        },
        ColorTheme.ALGERIA: {
            "primary": "006233",      # Vert
            "secondary": "D52B1E",    # Rouge
            "accent": "FFFFFF",       # Blanc
            "background": "FFFFFF",
            "text": "006233",
            "text_light": "FFFFFF"
        },
        ColorTheme.SWISS: {
            "primary": "FF0000",      # Rouge
            "secondary": "FFFFFF",    # Blanc
            "accent": "CC0000",       # Rouge foncé
            "background": "FFFFFF",
            "text": "333333",
            "text_light": "FFFFFF"
        },
        ColorTheme.TECH: {
            "primary": "0D1B2A",      # Bleu très foncé
            "secondary": "1B263B",    # Bleu foncé
            "accent": "00D9FF",       # Cyan
            "background": "0D1B2A",
            "text": "E0E1DD",
            "text_light": "FFFFFF"
        },
        ColorTheme.DARK: {
            "primary": "1A1A2E",
            "secondary": "16213E",
            "accent": "E94560",
            "background": "0F0F1A",
            "text": "EAEAEA",
            "text_light": "FFFFFF"
        }
    }
    
    # Templates de slides par type
    TEMPLATES = {
        PresentationType.PITCH_DECK: [
            {"layout": SlideLayout.TITLE, "title": "Titre"},
            {"layout": SlideLayout.BULLETS, "title": "Le Problème"},
            {"layout": SlideLayout.BULLETS, "title": "Notre Solution"},
            {"layout": SlideLayout.BULLETS, "title": "Proposition de Valeur"},
            {"layout": SlideLayout.CHART, "title": "Marché Cible"},
            {"layout": SlideLayout.BULLETS, "title": "Business Model"},
            {"layout": SlideLayout.CHART, "title": "Traction"},
            {"layout": SlideLayout.COMPARISON, "title": "Concurrence"},
            {"layout": SlideLayout.BULLETS, "title": "Équipe"},
            {"layout": SlideLayout.CHART, "title": "Projections Financières"},
            {"layout": SlideLayout.BULLETS, "title": "Demande d'Investissement"},
            {"layout": SlideLayout.THANK_YOU, "title": "Merci"}
        ],
        PresentationType.SALES_DECK: [
            {"layout": SlideLayout.TITLE, "title": "Titre"},
            {"layout": SlideLayout.SECTION, "title": "Qui Sommes-Nous"},
            {"layout": SlideLayout.BULLETS, "title": "Votre Challenge"},
            {"layout": SlideLayout.BULLETS, "title": "Notre Solution"},
            {"layout": SlideLayout.TWO_COLUMNS, "title": "Fonctionnalités Clés"},
            {"layout": SlideLayout.QUOTE, "title": "Témoignage Client"},
            {"layout": SlideLayout.COMPARISON, "title": "Avant / Après"},
            {"layout": SlideLayout.CHART, "title": "ROI"},
            {"layout": SlideLayout.BULLETS, "title": "Tarification"},
            {"layout": SlideLayout.BULLETS, "title": "Prochaines Étapes"},
            {"layout": SlideLayout.THANK_YOU, "title": "Questions?"}
        ],
        PresentationType.TRAINING: [
            {"layout": SlideLayout.TITLE, "title": "Formation"},
            {"layout": SlideLayout.BULLETS, "title": "Objectifs"},
            {"layout": SlideLayout.BULLETS, "title": "Agenda"},
            {"layout": SlideLayout.SECTION, "title": "Module 1"},
            {"layout": SlideLayout.TITLE_CONTENT, "title": "Concept 1"},
            {"layout": SlideLayout.TITLE_CONTENT, "title": "Exercice Pratique"},
            {"layout": SlideLayout.SECTION, "title": "Module 2"},
            {"layout": SlideLayout.BULLETS, "title": "Récapitulatif"},
            {"layout": SlideLayout.BULLETS, "title": "Ressources"},
            {"layout": SlideLayout.THANK_YOU, "title": "Merci"}
        ]
    }
    
    def __init__(self):
        self.output_dir = "outputs/presentations"
        os.makedirs(self.output_dir, exist_ok=True)
        self.generated_presentations: Dict[str, GeneratedPresentation] = {}
    
    async def generate(self, request: PresentationRequest) -> GeneratedPresentation:
        """Génère une présentation PPTX"""
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            raise HTTPException(
                status_code=500,
                detail="python-pptx non installé. Exécutez: pip install python-pptx"
            )
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
        
        theme = self.THEMES.get(request.theme, self.THEMES[ColorTheme.IA_FACTORY])
        
        slides_data = request.slides if request.slides else self._get_template_slides(request)
        
        slide_count = 0
        
        for slide_data in slides_data:
            self._add_slide(prs, slide_data, theme, request)
            slide_count += 1
        
        # Slide de contact si demandée
        if request.include_contact:
            contact_slide = SlideRequest(
                layout=SlideLayout.THANK_YOU,
                title="Contactez-nous",
                content={
                    "email": "contact@iafactory.ch",
                    "website": "www.iafactory.ch",
                    "linkedin": "linkedin.com/company/iafactory"
                }
            )
            self._add_slide(prs, contact_slide, theme, request)
            slide_count += 1
        
        # Sauvegarder
        pres_id = f"pres_{datetime.now().strftime('%Y%m%d%H%M%S')}_{request.presentation_type.value}"
        filename = f"{pres_id}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        
        prs.save(filepath)
        
        generated = GeneratedPresentation(
            id=pres_id,
            filename=filename,
            filepath=filepath,
            presentation_type=request.presentation_type,
            slide_count=slide_count,
            created_at=datetime.now(),
            download_url=f"/skills/pptx/download/{pres_id}"
        )
        
        self.generated_presentations[pres_id] = generated
        
        return generated
    
    def _add_slide(self, prs, slide_data: SlideRequest, theme: Dict, request: PresentationRequest):
        """Ajoute une slide à la présentation"""
        
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Choisir le layout approprié
        layout_index = self._get_layout_index(slide_data.layout)
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
        
        # Contenu selon le layout
        if slide_data.layout == SlideLayout.BULLETS and slide_data.content:
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                
                if isinstance(slide_data.content, list):
                    for i, item in enumerate(slide_data.content):
                        if i == 0:
                            tf.text = item
                        else:
                            p = tf.add_paragraph()
                            p.text = item
                            p.level = 0
        
        elif slide_data.layout == SlideLayout.TITLE_CONTENT and slide_data.content:
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                body.text = str(slide_data.content)
        
        elif slide_data.layout == SlideLayout.QUOTE and slide_data.content:
            # Ajouter une zone de texte pour la citation
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(11)
            height = Inches(3)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            if isinstance(slide_data.content, dict):
                quote_text = slide_data.content.get("quote", "")
                author = slide_data.content.get("author", "")
                
                p = tf.paragraphs[0]
                p.text = f'"{quote_text}"'
                p.font.size = Pt(28)
                p.font.italic = True
                p.alignment = PP_ALIGN.CENTER
                
                if author:
                    p2 = tf.add_paragraph()
                    p2.text = f"— {author}"
                    p2.font.size = Pt(18)
                    p2.alignment = PP_ALIGN.CENTER
        
        # Notes du présentateur
        if slide_data.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.notes
    
    def _get_layout_index(self, layout: SlideLayout) -> int:
        """Retourne l'index du layout PowerPoint"""
        layout_map = {
            SlideLayout.TITLE: 0,
            SlideLayout.TITLE_CONTENT: 1,
            SlideLayout.SECTION: 2,
            SlideLayout.TWO_COLUMNS: 3,
            SlideLayout.COMPARISON: 4,
            SlideLayout.BULLETS: 1,
            SlideLayout.QUOTE: 5,
            SlideLayout.IMAGE_FULL: 6,
            SlideLayout.CHART: 1,
            SlideLayout.THANK_YOU: 0
        }
        return layout_map.get(layout, 1)
    
    def _get_template_slides(self, request: PresentationRequest) -> List[SlideRequest]:
        """Génère les slides à partir du template"""
        template = self.TEMPLATES.get(request.presentation_type, [])
        
        slides = []
        for slide_template in template:
            slides.append(SlideRequest(
                layout=slide_template.get("layout", SlideLayout.TITLE_CONTENT),
                title=slide_template.get("title", ""),
                content=None,
                notes=None
            ))
        
        # Personnaliser la première slide
        if slides:
            slides[0].title = request.title
            slides[0].subtitle = request.subtitle
        
        return slides
    
    def get_template(self, presentation_type: PresentationType) -> List[Dict]:
        """Retourne le template pour un type de présentation"""
        return self.TEMPLATES.get(presentation_type, [])
    
    def list_presentations(self) -> List[GeneratedPresentation]:
        """Liste toutes les présentations générées"""
        return list(self.generated_presentations.values())


# Instance globale
pptx_generator = PptxGenerator()


# Routes API

@router.post("/generate", response_model=Dict[str, Any])
async def generate_presentation(request: PresentationRequest):
    """
    Génère une présentation PPTX professionnelle
    
    Types supportés:
    - pitch_deck: Pitch investisseurs
    - sales_deck: Présentation commerciale
    - training: Formation
    - product_demo: Démo produit
    """
    pres = await pptx_generator.generate(request)
    return {
        "status": "success",
        "presentation_id": pres.id,
        "filename": pres.filename,
        "slide_count": pres.slide_count,
        "download_url": pres.download_url
    }


@router.get("/templates/{presentation_type}")
async def get_template(presentation_type: PresentationType):
    """Retourne le template pour un type de présentation"""
    template = pptx_generator.get_template(presentation_type)
    return {
        "presentation_type": presentation_type.value,
        "slides": template,
        "slide_count": len(template)
    }


@router.get("/templates")
async def list_templates():
    """Liste tous les templates disponibles"""
    return {
        pres_type.value: {
            "slides": template,
            "slide_count": len(template)
        }
        for pres_type, template in pptx_generator.TEMPLATES.items()
    }


@router.get("/themes")
async def list_themes():
    """Liste tous les thèmes de couleurs disponibles"""
    return {
        theme.value: colors
        for theme, colors in pptx_generator.THEMES.items()
    }


@router.get("/presentations")
async def list_presentations():
    """Liste toutes les présentations générées"""
    return pptx_generator.list_presentations()


@router.get("/download/{pres_id}")
async def download_presentation(pres_id: str):
    """Télécharge une présentation générée"""
    from fastapi.responses import FileResponse
    
    pres = pptx_generator.generated_presentations.get(pres_id)
    if not pres:
        raise HTTPException(status_code=404, detail="Presentation not found")
    
    return FileResponse(
        path=pres.filepath,
        filename=pres.filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
